"""typeset 工具测试（typeset_generate）。

用 PHXSC_WORKDIR 环境变量把沙箱 workdir 指到 tmp_path 内目录，不触碰真实
workspace/。覆盖：简单笔记生成 PPTX（页数/标题/文件有效性）、二级标题分页与
页眉、列表项、表格（表头加粗/行列数）、超长自动分页、source="plans"、
源不存在与路径逃逸的结构化错误、原子保存（无 .tmp 残留）、工具注册与
typeset 模式工具集。测完清理。
"""

import os
import shutil

import pytest
from pptx import Presentation
from pptx.util import Pt

from phxsc.agent.tools import Tool, ToolRegistry
from phxsc.tools import notes as notes_tools
from phxsc.tools import typeset as typeset_tools


@pytest.fixture
def typeset_env(tmp_path, monkeypatch):
    """把 PHXSC_WORKDIR 指到 tmp_path 内目录，测完清理。"""
    workdir = tmp_path / "work"
    workdir.mkdir()
    monkeypatch.setenv("PHXSC_WORKDIR", str(workdir))
    yield workdir
    shutil.rmtree(workdir, ignore_errors=True)


def _write_source(workdir, subdir, name, content):
    """向 workdir/<subdir>/<name> 写入 Markdown 源文件。"""
    path = workdir / subdir / name
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(content, encoding="utf-8")
    return path


def _open_generated(workdir):
    """打开 workdir/typeset/ 下唯一的 pptx。"""
    files = list((workdir / "typeset").glob("*.pptx"))
    assert files, "typeset 目录下没有生成 .pptx"
    return Presentation(str(files[0]))


def _all_text(prs):
    """收集整份演示文稿所有可见文本（含表格单元格）。"""
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(r.text for r in para.runs).strip()
                    if t:
                        texts.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        t = cell.text.strip()
                        if t:
                            texts.append(t)
    return texts


def _header_of(slide):
    """返回幻灯片中 top 最小的文本形状首行文本（内容页页眉）。"""
    best, best_top = None, None
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        texts = [t for t in ("".join(r.text for r in p.runs).strip()
                             for p in shape.text_frame.paragraphs) if t]
        if not texts:
            continue
        if best_top is None or shape.top < best_top:
            best_top, best = shape.top, texts[0]
    return best


def _content_slide(prs, needle):
    """找到包含 needle 文本的内容页（即该页页眉为章节名的页面）。"""
    for slide in prs.slides:
        texts = _all_text_for_slide(slide)
        if needle in texts:
            return slide
    return None


def _all_text_for_slide(slide):
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = "".join(r.text for r in para.runs).strip()
                if t:
                    texts.append(t)
    return texts


def _shape_fills(prs):
    """收集所有 shape 的填充色（hex 字符串），容忍无/继承填充。"""
    colors = []
    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                colors.append(str(shape.fill.fore_color.rgb))
            except Exception:
                continue
    return colors


def _run_colors_by_text(prs):
    """文本 → run 颜色列表（hex 字符串），容忍 None/继承色。"""
    mapping = {}
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        color = str(run.font.color.rgb)
                    except Exception:
                        continue
                    mapping.setdefault(run.text, []).append(color)
    return mapping


class TestSimpleNote:
    def test_generates_pptx_with_cover_and_content(self, typeset_env):
        _write_source(typeset_env, "notes", "总结.md", "# 钙钛矿综述\n\n第一段内容。\n\n第二段内容。\n")
        out = typeset_tools.typeset_generate.fn(title="总结")
        assert out.startswith("已生成 typeset/总结.pptx")
        assert "页，" in out and "字符源内容" in out

        pptx_path = typeset_env / "typeset" / "总结.pptx"
        assert pptx_path.is_file()
        assert pptx_path.stat().st_size > 0

        prs = Presentation(str(pptx_path))
        assert len(prs.slides) >= 2
        assert "钙钛矿综述" in _all_text(prs)
        assert "第一段内容。" in _all_text(prs)

    def test_cover_subtitle_mentions_source_and_update(self, typeset_env):
        _write_source(typeset_env, "notes", "n.md", "# 标题\n\n内容。\n")
        typeset_tools.typeset_generate.fn(title="n")
        prs = _open_generated(typeset_env)
        cover_texts = _all_text_for_slide(prs.slides[0])
        subtitle = next((t for t in cover_texts if "来源" in t), "")
        assert "notes/n.md" in subtitle
        assert "更新于" in subtitle


class TestSectionPagination:
    def test_sections_split_across_slides_with_header(self, typeset_env):
        _write_source(
            typeset_env, "notes", "分章.md",
            "# 总标题\n\n"
            "## 第一章\n\n第一章内容甲。\n\n"
            "## 第二章\n\n第二章内容乙。\n",
        )
        typeset_tools.typeset_generate.fn(title="分章")
        prs = _open_generated(typeset_env)
        texts = _all_text(prs)
        assert "第一章内容甲。" in texts and "第二章内容乙。" in texts

        slide_a = _content_slide(prs, "第一章内容甲。")
        slide_b = _content_slide(prs, "第二章内容乙。")
        assert slide_a is not None and slide_b is not None
        assert slide_a is not slide_b
        assert _header_of(slide_a) == "第一章"
        assert _header_of(slide_b) == "第二章"

    def test_header_font_size_is_12pt(self, typeset_env):
        _write_source(typeset_env, "notes", "h.md", "# T\n\n## 章节\n\n内容。\n")
        typeset_tools.typeset_generate.fn(title="h")
        prs = _open_generated(typeset_env)
        slide = _content_slide(prs, "内容。")
        assert slide is not None
        header_shape = min(
            (s for s in slide.shapes if s.has_text_frame
             and "".join(r.text for r in s.text_frame.paragraphs[0].runs).strip()),
            key=lambda s: s.top,
        )
        run = header_shape.text_frame.paragraphs[0].runs[0]
        assert run.font.size == Pt(12)


class TestLists:
    def test_all_list_items_preserved(self, typeset_env):
        items = [f"要点{i}" for i in range(1, 11)]
        note = "# 清单\n\n" + "\n".join(f"- {it}" for it in items) + "\n"
        _write_source(typeset_env, "notes", "清单.md", note)
        typeset_tools.typeset_generate.fn(title="清单")
        prs = _open_generated(typeset_env)
        texts = _all_text(prs)
        for it in items:
            assert any(t.endswith(it) for t in texts), f"列表项 {it} 丢失"


class TestTable:
    def test_table_rows_cols_and_bold_header(self, typeset_env):
        _write_source(
            typeset_env, "notes", "表.md",
            "# 对比\n\n| 方法 | 效果 |\n|------|------|\n| A | 好 |\n| B | 差 |\n",
        )
        typeset_tools.typeset_generate.fn(title="表")
        prs = _open_generated(typeset_env)
        table = next(shape.table for slide in prs.slides
                     for shape in slide.shapes if shape.has_table)
        assert len(table.rows) == 3
        assert len(table.columns) == 2
        assert table.cell(0, 0).text.strip() == "方法"
        assert table.cell(2, 1).text.strip() == "差"
        header_run = table.cell(0, 0).text_frame.paragraphs[0].runs[0]
        body_run = table.cell(1, 0).text_frame.paragraphs[0].runs[0]
        assert header_run.font.bold is True
        assert body_run.font.bold is not True


class TestLongContentPagination:
    def test_many_paragraphs_paginate_beyond_three_slides(self, typeset_env):
        note = "# 长文\n\n" + "\n\n".join(f"这是第{i}段内容。" for i in range(22)) + "\n"
        _write_source(typeset_env, "notes", "长文.md", note)
        typeset_tools.typeset_generate.fn(title="长文")
        prs = _open_generated(typeset_env)
        assert len(prs.slides) > 3


class TestCodeSlide:
    """P3-5：代码块按 MAX_CHARS 截断，超长不再整块渲一页溢出。"""

    MARKER = "…（代码过长已截断）"

    def test_long_code_truncated_with_marker(self, typeset_env):
        code = "x" * (typeset_tools.MAX_CHARS + 500)
        _write_source(typeset_env, "notes", "longcode.md",
                      "# 代码\n\n```python\n" + code + "\n```\n")
        typeset_tools.typeset_generate.fn(title="longcode")
        prs = _open_generated(typeset_env)
        texts = _all_text(prs)
        rendered = next(t for t in texts if self.MARKER in t)
        body = rendered.split("\n" + self.MARKER)[0]
        assert len(body) == typeset_tools.MAX_CHARS
        assert "已截断" in rendered

    def test_short_code_not_truncated(self, typeset_env):
        _write_source(typeset_env, "notes", "shortcode.md",
                      "# 代码\n\n```python\nx = 1\n```\n")
        typeset_tools.typeset_generate.fn(title="shortcode")
        prs = _open_generated(typeset_env)
        texts = _all_text(prs)
        assert "x = 1" in texts
        assert not any(self.MARKER in t for t in texts)


class TestSourcePlans:
    def test_reads_from_plans_directory(self, typeset_env):
        _write_source(typeset_env, "plans", "plan.md", "# 研究计划\n\n内容。\n")
        out = typeset_tools.typeset_generate.fn(title="plan", source="plans")
        assert out.startswith("已生成 typeset/plan.pptx")
        prs = _open_generated(typeset_env)
        assert "研究计划" in _all_text(prs)

    def test_unknown_source_rejected(self, typeset_env):
        _write_source(typeset_env, "notes", "n.md", "内容")
        result = typeset_tools.typeset_generate.fn(title="n", source="nonexistent")
        assert isinstance(result, dict)
        assert set(result) == {"error", "reason", "fix_hint"}


class TestStructuredErrors:
    def test_missing_source_returns_structured_error(self, typeset_env):
        result = typeset_tools.typeset_generate.fn(title="不存在的笔记")
        assert isinstance(result, dict)
        assert set(result) == {"error", "reason", "fix_hint"}
        assert "不存在" in result["error"]

    def test_path_escape_cleaned_then_rejected(self, typeset_env):
        result = typeset_tools.typeset_generate.fn(title="../../etc/passwd")
        assert isinstance(result, dict)
        assert set(result) == {"error", "reason", "fix_hint"}

    def test_escape_after_cleaning_rejected(self, typeset_env, monkeypatch):
        monkeypatch.setattr(typeset_tools, "_clean_title", lambda t: "../../evil.md")
        result = typeset_tools.typeset_generate.fn(title="x")
        assert isinstance(result, dict)
        assert set(result) == {"error", "reason", "fix_hint"}


class TestAtomicSave:
    def test_no_tmp_leftover_after_generation(self, typeset_env):
        _write_source(typeset_env, "notes", "a.md", "# A\n\n内容。\n")
        typeset_tools.typeset_generate.fn(title="a")
        assert list((typeset_env / "typeset").glob("*.tmp")) == []
        assert list((typeset_env / "typeset").glob("*.pptx")) != []


class TestThemes:
    """Q10 PPT 审美升级：多套配色主题 + style 参数路由。"""

    def test_theme_dicts_have_required_keys(self):
        for name in ("academic", "deep", "warm"):
            assert set(typeset_tools.THEMES[name]) >= {
                "primary", "secondary", "accent", "background", "text",
            }
        assert typeset_tools.THEMES["academic"]["primary"] == "1F4E79"
        assert typeset_tools.THEMES["deep"]["background"] == "1E242B"
        assert typeset_tools.THEMES["warm"]["primary"] == "B85042"

    @pytest.mark.parametrize("style", ["academic", "deep", "warm"])
    def test_style_uses_theme_primary_fill(self, typeset_env, style):
        _write_source(typeset_env, "notes", "n.md", "# 标题A\n\n## 章节\n\n内容。\n")
        out = typeset_tools.typeset_generate.fn(title="n", style=style)
        assert f"style={style}" in out
        prs = _open_generated(typeset_env)
        assert typeset_tools.THEMES[style]["primary"] in _shape_fills(prs)

    @pytest.mark.parametrize("style", ["academic", "deep", "warm"])
    def test_cover_title_white_and_section_title_primary(self, typeset_env, style):
        _write_source(typeset_env, "notes", "n.md", "# 我的标题\n\n## 我的章节\n\n内容。\n")
        typeset_tools.typeset_generate.fn(title="n", style=style)
        prs = _open_generated(typeset_env)
        colors = _run_colors_by_text(prs)
        assert colors.get("我的标题"), "封面标题缺失"
        assert all(c == "FFFFFF" for c in colors["我的标题"])
        assert colors.get("我的章节"), "章节标题缺失"
        assert all(c == typeset_tools.THEMES[style]["primary"] for c in colors["我的章节"])

    def test_table_header_uses_primary_background_white_text(self, typeset_env):
        _write_source(
            typeset_env, "notes", "t.md",
            "# 表\n\n| A | B |\n|---|---|\n| 1 | 2 |\n",
        )
        typeset_tools.typeset_generate.fn(title="t", style="warm")
        prs = _open_generated(typeset_env)
        table = next(shape.table for slide in prs.slides
                     for shape in slide.shapes if shape.has_table)
        assert str(table.cell(0, 0).fill.fore_color.rgb) == typeset_tools.THEMES["warm"]["primary"]
        header_run = table.cell(0, 0).text_frame.paragraphs[0].runs[0]
        assert str(header_run.font.color.rgb) == "FFFFFF"

    def test_code_block_background_from_theme(self, typeset_env):
        _write_source(typeset_env, "notes", "c.md", "# 代码\n\n```python\nx = 1\n```\n")
        typeset_tools.typeset_generate.fn(title="c", style="academic")
        prs = _open_generated(typeset_env)
        assert typeset_tools.THEMES["academic"]["code_bg"] in _shape_fills(prs)

    def test_auto_routes_review_to_academic(self, typeset_env):
        _write_source(typeset_env, "notes", "rev.md", "# 钙钛矿综述\n\n内容。\n")
        out = typeset_tools.typeset_generate.fn(title="rev")
        assert "style=academic" in out

    def test_auto_routes_report_to_deep(self, typeset_env):
        _write_source(typeset_env, "notes", "rep.md", "# 年度报告\n\n内容。\n")
        out = typeset_tools.typeset_generate.fn(title="rep")
        assert "style=deep" in out

    def test_auto_routes_english_survey_to_academic(self, typeset_env):
        _write_source(typeset_env, "notes", "en.md", "# A Review\n\n内容。\n")
        out = typeset_tools.typeset_generate.fn(title="en")
        assert "style=academic" in out

    def test_auto_defaults_to_academic(self, typeset_env):
        _write_source(typeset_env, "notes", "other.md", "# 随便什么内容\n\n正文。\n")
        out = typeset_tools.typeset_generate.fn(title="other")
        assert "style=academic" in out

    def test_invalid_style_returns_structured_error(self, typeset_env):
        _write_source(typeset_env, "notes", "n.md", "# 标题\n\n内容。\n")
        result = typeset_tools.typeset_generate.fn(title="n", style="neon")
        assert isinstance(result, dict)
        assert set(result) == {"error", "reason", "fix_hint"}
        assert "style" in result["error"]

    def test_style_param_appears_in_schema(self):
        props = typeset_tools.typeset_generate.parameters["properties"]
        assert props["style"]["default"] == "auto"
        assert typeset_tools.typeset_generate.parameters["required"] == ["title"]


class TestToolRegistration:
    def test_decorated_as_tool(self):
        assert isinstance(typeset_tools.typeset_generate, Tool)
        assert typeset_tools.typeset_generate.name == "typeset_generate"
        assert typeset_tools.typeset_generate.mode == {"typeset"}
        assert typeset_tools.typeset_generate.parameters["required"] == ["title"]
        assert typeset_tools.typeset_generate.parameters["properties"]["source"]["default"] == "notes"

    def test_typeset_mode_toolset(self):
        reg = ToolRegistry()
        reg.register_all([
            typeset_tools.typeset_generate,
            notes_tools.notes_read,
            notes_tools.notes_list,
        ])
        names = {t["function"]["name"] for t in reg.get_tools("typeset")}
        assert {"typeset_generate", "notes_read", "notes_list"} <= names
        plan_names = {t["function"]["name"] for t in reg.get_tools("plan")}
        assert "typeset_generate" not in plan_names
