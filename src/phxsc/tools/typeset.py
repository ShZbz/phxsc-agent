"""沙箱文档生成工具：typeset_generate。

把 notes/ 或 plans/ 下的 Markdown 笔记转换成 PPTX 演示文稿，输出到
<workdir>/typeset/。读取路径过沙箱白名单（safe_read_path），逃逸拒绝并返回
{error, reason, fix_hint} 结构化错误。v0.0.1 只做 PPTX 基础版（DOCX/PDF 后置）。

Markdown 支持（基础版）：
- # 一级标题 → 封面页（标题 + 副标题：来源与更新时间）
- ## 二级标题 → 新章节标题页，后续内容页以章节名为页眉
- ### 三级标题 → 加粗段落
- 正文段落 / - 列表 / > 引用 / 图片行 → 内容页（容量内自动分页）
- 表格（| 分隔）→ python-pptx 表格，首行表头加粗
- ``` 代码块 → 等宽字体文本框（独立页）
- 行内链接 [text](url) 保留文字；行内图片标注 [图片略]

排版：16:9 宽屏，微软雅黑（中）+ Calibri（西），多套配色主题（THEMES），
style 参数控制：auto 按源内容关键词猜（综述/Review/survey → academic，
报告/Report → deep，其他 → academic），也可显式指定 academic/deep/warm。
封面标题白字 40pt，章节页左侧 primary 色条，内容页页眉 primary 色 12pt，
正文 theme["text"]，表格表头 primary 底白字，代码块 code_bg 底；页眉 12pt；
边距 ≥0.5in；行距 1.3；每页约 6 段落块或 8 列表项。
"""

import os
import re
import subprocess
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from phxsc.agent.tools import tool
from phxsc.sandbox.paths import safe_read_path, safe_write_path
from phxsc.tools.notes import _clean_title

try:
    import pymupdf as _fitz
except ImportError:  # PyMuPDF 旧版本仅提供 fitz 顶层包
    import fitz as _fitz

TYPESET_DIR = "typeset"

# 多套配色主题（Q10 PPT 审美升级）：所有渲染颜色从 theme 取，不用硬编码。
# academic=学术蓝（默认）；deep=深色科技；warm=暖色简约。
THEMES = {
    "academic": {
        "primary": "1F4E79",     # 主色：标题 / 章节条 / 表头 / 页眉
        "secondary": "E8EEF7",   # 辅色：表格隔行 / 引用（浅蓝灰，更克制）
        "accent": "1F4E79",      # 强调色：装饰条（与 primary 同色，单 accent 锁定）
        "background": "FFFFFF",  # 背景
        "text": "333333",        # 正文深灰
        "code_bg": "F1F5F9",     # 代码块底
    },
    "deep": {
        "primary": "0D9488",
        "secondary": "2A3440",
        "accent": "0D9488",
        "background": "1E242B",  # 微暖深灰，避免纯 GitHub 深蓝黑
        "text": "F1F5F9",
        "code_bg": "1A2333",     # 深色底，与页面一致
    },
    "warm": {
        "primary": "B85042",
        "secondary": "E7E8D1",
        "accent": "B85042",      # 与 primary 同色，单 accent 锁定
        "background": "FAF6F0",
        "text": "333333",
        "code_bg": "F4EFE6",     # 暖浅底，与暖主题一致
    },
}
STYLE_NAMES = tuple(THEMES)

ACCENT = "1F4E79"       # 学术蓝：标题 / 强调 / 表头（默认主题主色）
BODY = "333333"         # 正文深灰
HEADER_GRAY = "808080"  # 页眉灰
QUOTE = "4F6E8F"        # 引用蓝灰
CODE_BG = "F1F5F9"      # 代码块浅灰底
CODE_BORDER = "D9D9D9"

MAX_BLOCKS = 6      # 每内容页最多段落/引用块
MAX_LIST = 8        # 每内容页最多列表项
MAX_CHARS = 1500    # 每内容页估算字符上限（防溢出）

TITLE_SIZE = 40     # 封面标题
SECTION_SIZE = 36   # 章节页标题
BODY_SIZE = 15      # 正文
HEADER_SIZE = 12    # 页眉


def _workdir() -> str:
    """workdir：PHXSC_WORKDIR 环境变量优先，默认 <项目根>/workspace；确保 typeset/ 存在。"""
    env = os.environ.get("PHXSC_WORKDIR")
    if env:
        workdir = env
    else:
        workdir = str(Path(__file__).resolve().parents[3] / "workspace")
    os.makedirs(os.path.join(workdir, TYPESET_DIR), exist_ok=True)
    return workdir


def _err(error: str, reason: str, fix_hint: str) -> dict:
    """结构化错误 dict。"""
    return {"error": error, "reason": reason, "fix_hint": fix_hint}


def _denied_to_err(exc: ValueError) -> dict:
    """把 safe_*_path 的 ValueError（内含 reason/fix_hint）解析为错误 dict。"""
    msg = str(exc)
    reason = "ValueError"
    fix_hint = "使用 workdir 内的路径"
    for seg in msg.split("|"):
        seg = seg.strip()
        if seg.startswith("reason:"):
            reason = seg[len("reason:") :].strip()
        elif seg.startswith("fix_hint:"):
            fix_hint = seg[len("fix_hint:") :].strip()
    return _err(f"路径校验失败：{msg}", reason, fix_hint)


def _resolve_theme(style: str, content: str) -> str | None:
    """把 style 参数解析为 THEMES 键；auto 按源内容关键词路由，非法值返回 None。

    auto 规则：含 综述/Review/survey → academic；含 报告/Report → deep；
    其他 → academic。显式 academic/deep/warm 原样返回。
    """
    if style in THEMES:
        return style
    if style == "auto":
        low = content.lower()
        if any(k in low for k in ("综述", "review", "survey")):
            return "academic"
        if any(k in low for k in ("报告", "report")):
            return "deep"
        return "academic"
    return None


def _clean_inline(text: str) -> str:
    """行内 Markdown 基础清理：链接保留文字，图片标注，剥掉强调/行内代码记号。"""
    text = re.sub(r"!\[([^\]]*)\]\([^)]*\)", "[图片略]", text)
    text = re.sub(r"\[([^\]]*)\]\([^)]*\)", r"\1", text)
    text = text.replace("**", "").replace("__", "").replace("`", "")
    text = re.sub(r"\*([^*]+)\*", r"\1", text)
    return text.strip()


def _parse_table(rows: list[str]) -> list[list[str]] | None:
    """把 markdown 表格行解析为二维单元格列表；去掉分隔行，非表格返回 None。"""
    table = []
    for raw in rows:
        cells = [c.strip() for c in raw.strip().strip("|").split("|")]
        if all(re.fullmatch(r":?-{1,}:?", c) or c == "" for c in cells):
            continue  # 分隔行 |---|
        table.append(cells)
    if len(table) >= 2:
        return table
    return None


def _parse_blocks(content: str) -> list[dict]:
    """把 Markdown 拆成顶层块序列（heading / para / list / quote / image / table / code）。"""
    blocks: list[dict] = []
    lines = content.splitlines()
    i, n = 0, len(lines)
    while i < n:
        raw = lines[i].rstrip()
        stripped = raw.strip()
        if not stripped:
            i += 1
            continue
        if stripped.startswith("```"):
            code = []
            i += 1
            while i < n and not lines[i].strip().startswith("```"):
                code.append(lines[i])
                i += 1
            i += 1  # 跳过闭合 fence
            blocks.append({"type": "code", "text": "\n".join(code)})
            continue
        m_head = re.match(r"^(#{1,6})\s+(.*)$", stripped)
        if m_head:
            level = len(m_head.group(1))
            blocks.append({"type": "heading", "level": level, "text": _clean_inline(m_head.group(2))})
            i += 1
            continue
        if stripped.startswith("|"):
            rows = [stripped]
            i += 1
            while i < n and lines[i].strip().startswith("|"):
                rows.append(lines[i].strip())
                i += 1
            parsed = _parse_table(rows)
            if parsed:
                blocks.append({"type": "table", "rows": parsed})
                continue
            blocks.append({"type": "para", "text": stripped})
            continue
        if stripped.startswith(">"):
            blocks.append({"type": "quote", "text": _clean_inline(stripped.lstrip(">").strip())})
            i += 1
            continue
        if re.match(r"^!\[[^\]]*\]\([^)]*\)\s*$", stripped):
            blocks.append({"type": "image", "text": "图片"})
            i += 1
            continue
        if re.match(r"^[-*+]\s+", stripped):
            blocks.append({"type": "list", "text": _clean_inline(re.sub(r"^[-*+]\s+", "", stripped))})
            i += 1
            continue
        if re.match(r"^\d+[.)]\s+", stripped):
            blocks.append({"type": "list", "text": _clean_inline(re.sub(r"^\d+[.)]\s+", "", stripped))})
            i += 1
            continue
        blocks.append({"type": "para", "text": _clean_inline(stripped)})
        i += 1
    return blocks


def _build_slides(content: str, source: str, fname: str, target_src: str) -> list[dict]:
    """把解析后的块组织成幻灯片规格列表（cover/section/content/table/code）。"""
    blocks = _parse_blocks(content)
    slides: list[dict] = []
    cover = None
    current_header = None
    content_buf: dict | None = None

    def flush() -> None:
        nonlocal content_buf
        if content_buf and content_buf["blocks"]:
            slides.append({
                "kind": "content",
                "header": content_buf["header"],
                "blocks": content_buf["blocks"],
            })
        content_buf = None

    def start_content(header: str) -> None:
        nonlocal content_buf
        content_buf = {"header": header, "blocks": [], "chars": 0, "list_count": 0}

    def add_block(block: dict, header: str) -> None:
        nonlocal content_buf
        if content_buf is None:
            start_content(header)
        kind = block["type"]
        text = block.get("text", "")
        para_used = sum(b["type"] in ("para", "quote") for b in content_buf["blocks"])
        over_budget = (
            (kind in ("para", "quote") and para_used >= MAX_BLOCKS)
            or (kind == "list" and content_buf["list_count"] >= MAX_LIST)
            or content_buf["chars"] + len(text) + 8 > MAX_CHARS
        )
        if over_budget and content_buf["blocks"]:
            flush()
            start_content(header)
        content_buf["blocks"].append(block)
        content_buf["chars"] += len(text) + 8
        if kind == "list":
            content_buf["list_count"] += 1

    for b in blocks:
        t = b["type"]
        if t == "heading" and b["level"] == 1:
            if cover is None:
                cover = b["text"]
                current_header = b["text"]
            else:
                flush()
                slides.append({"kind": "section", "title": b["text"]})
                current_header = b["text"]
            continue
        if t == "heading" and b["level"] == 2:
            flush()
            slides.append({"kind": "section", "title": b["text"]})
            current_header = b["text"]
            continue
        if t == "heading":  # ### 三级标题 → 内容页加粗段
            flush()
            start_content(current_header or cover or "内容")
            content_buf["blocks"].append({"type": "para", "text": b["text"], "strong": True})
            content_buf["chars"] += len(b["text"]) + 8
            continue
        if t == "table":
            flush()
            slides.append({"kind": "table", "header": current_header or cover or "内容", "rows": b["rows"]})
            continue
        if t == "code":
            flush()
            slides.append({"kind": "code", "header": current_header or cover or "内容", "text": b["text"]})
            continue
        add_block(b, current_header or cover or "内容")
    flush()

    if cover is None:
        cover = fname[:-3] if fname.endswith(".md") else fname
        subtitle = f"来源：{source}/{fname}"
    else:
        mtime = datetime.fromtimestamp(os.path.getmtime(target_src)).strftime("%Y-%m-%d %H:%M")
        subtitle = f"来源：{source}/{fname} · 更新于 {mtime}"
    slides.insert(0, {"kind": "cover", "title": cover, "subtitle": subtitle})
    return slides


def _set_font(run, *, latin: str = "Calibri", ea: str = "微软雅黑", size: int | None = None,
              bold: bool | None = None, color: str | None = None) -> None:
    """设置 run 字体：latin 走 Calibri/西文，ea 走微软雅黑（东亚字符）。"""
    f = run.font
    f.name = latin
    if size is not None:
        f.size = Pt(size)
    if bold is not None:
        f.bold = bold
    if color is not None:
        f.color.rgb = RGBColor.from_string(color)
    rPr = run._r.get_or_add_rPr()
    latin_el = rPr.find(qn("a:latin"))
    ea_el = rPr.find(qn("a:ea"))
    if ea_el is None:
        ea_el = rPr.makeelement(qn("a:ea"), {})
        if latin_el is not None:
            latin_el.addnext(ea_el)
        else:
            rPr.append(ea_el)
    ea_el.set("typeface", ea)


def _add_header(slide, text: str, theme: dict) -> None:
    """内容页页眉：章节名 12pt primary 色 + 强调色下划线。"""
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.22), Inches(12.1), Inches(0.4))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    _set_font(run, size=HEADER_SIZE, color=theme["primary"])
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.62), Inches(12.1), Inches(0.03))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor.from_string(theme["accent"])
    bar.line.fill.background()


def _add_cover(prs, blank, spec: dict, theme: dict) -> None:
    slide = prs.slides.add_slide(blank)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor.from_string(theme["background"])
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.35), Inches(13.333), Inches(1.35))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor.from_string(theme["primary"])
    bar.line.fill.background()
    box = slide.shapes.add_textbox(Inches(0.9), Inches(2.6), Inches(11.5), Inches(1.0))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = spec["title"]
    _set_font(run, size=TITLE_SIZE, bold=True, color="FFFFFF")
    p2 = tf.add_paragraph()
    p2.space_before = Pt(12)
    run2 = p2.add_run()
    run2.text = spec["subtitle"]
    _set_font(run2, size=16, color=HEADER_GRAY)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(4.05), Inches(2.2), Inches(0.05))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor.from_string(theme["accent"])
    accent.line.fill.background()


def _add_section(prs, blank, spec: dict, theme: dict) -> None:
    slide = prs.slides.add_slide(blank)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor.from_string(theme["background"])
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(2.85), Inches(0.12), Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor.from_string(theme["primary"])
    bar.line.fill.background()
    box = slide.shapes.add_textbox(Inches(1.25), Inches(2.85), Inches(11.2), Inches(1.6))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = spec["title"]
    _set_font(run, size=SECTION_SIZE, bold=True, color=theme["primary"])


def _add_content(prs, blank, spec: dict, theme: dict) -> None:
    slide = prs.slides.add_slide(blank)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor.from_string(theme["background"])
    _add_header(slide, spec["header"], theme)
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.95), Inches(12.1), Inches(6.0))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for b in spec["blocks"]:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.line_spacing = 1.3
        if b["type"] == "list":
            run = p.add_run()
            run.text = "• " + b["text"]
            _set_font(run, size=BODY_SIZE, color=theme["text"])
            p.space_after = Pt(6)
        elif b["type"] == "quote":
            run = p.add_run()
            run.text = b["text"]
            _set_font(run, size=BODY_SIZE, color=theme["primary"])
            p.level = 1
            p.space_after = Pt(6)
        elif b["type"] == "image":
            run = p.add_run()
            run.text = "[图片略]"
            _set_font(run, size=13, color=HEADER_GRAY)
            p.space_after = Pt(6)
        elif b.get("strong"):  # ### 三级标题 → 20pt primary 色加粗
            run = p.add_run()
            run.text = b["text"]
            _set_font(run, size=20, bold=True, color=theme["primary"])
            p.space_after = Pt(8)
        else:  # para
            run = p.add_run()
            run.text = b["text"]
            _set_font(run, size=BODY_SIZE, color=theme["text"])
            p.space_after = Pt(8)


def _fill_cell(cell, text: str, *, bold: bool, color: str, size: int) -> None:
    cell.margin_left = Inches(0.12)
    cell.margin_right = Inches(0.12)
    cell.margin_top = Inches(0.05)
    cell.margin_bottom = Inches(0.05)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    _set_font(run, size=size, bold=bold, color=color)


def _add_table_slide(prs, blank, spec: dict, theme: dict) -> None:
    slide = prs.slides.add_slide(blank)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor.from_string(theme["background"])
    _add_header(slide, spec["header"], theme)
    rows = spec["rows"]
    nrows, ncols = len(rows), len(rows[0])
    table_shape = slide.shapes.add_table(nrows, ncols, Inches(0.6), Inches(1.2), Inches(12.1), Inches(2.0))
    table = table_shape.table
    lengths = [0.0] * ncols
    for r in rows:
        for c in range(min(ncols, len(r))):
            lengths[c] = max(lengths[c], len(r[c]))
    total = sum(lengths) or ncols
    for c in range(ncols):
        table.columns[c].width = Inches(max(1.4, 11.9 * lengths[c] / total))
    for r in range(nrows):
        table.rows[r].height = Inches(0.5)
    for c in range(ncols):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor.from_string(theme["primary"])
        _fill_cell(cell, rows[0][c] if c < len(rows[0]) else "", bold=True, color="FFFFFF", size=13)
    for r in range(1, nrows):
        for c in range(ncols):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor.from_string(
                theme["secondary"] if r % 2 == 0 else theme["background"]
            )
            _fill_cell(cell, rows[r][c] if c < len(rows[r]) else "", bold=False, color=theme["text"], size=12)


def _add_code_slide(prs, blank, spec: dict, theme: dict) -> None:
    slide = prs.slides.add_slide(blank)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor.from_string(theme["background"])
    _add_header(slide, spec["header"], theme)
    box = slide.shapes.add_textbox(Inches(0.6), Inches(1.0), Inches(12.1), Inches(6.0))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor.from_string(theme["code_bg"])
    box.line.color.rgb = RGBColor.from_string(CODE_BORDER)
    box.line.width = Pt(1)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.2)
    tf.margin_bottom = Inches(0.2)
    p = tf.paragraphs[0]
    run = p.add_run()
    code_text = spec["text"] or " "
    if len(code_text) > MAX_CHARS:
        code_text = code_text[:MAX_CHARS] + "\n…（代码过长已截断）"
    run.text = code_text
    _set_font(run, latin="Consolas", ea="微软雅黑", size=12, color=theme["text"])


def _save_pptx(slides: list[dict], out_path: str, theme: dict) -> None:
    """渲染并原子保存：先写 .tmp 再 os.replace（.tmp 与目标同目录）。"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for spec in slides:
        kind = spec["kind"]
        if kind == "cover":
            _add_cover(prs, blank, spec, theme)
        elif kind == "section":
            _add_section(prs, blank, spec, theme)
        elif kind == "content":
            _add_content(prs, blank, spec, theme)
        elif kind == "table":
            _add_table_slide(prs, blank, spec, theme)
        else:
            _add_code_slide(prs, blank, spec, theme)
    tmp = out_path + ".tmp"
    prs.save(tmp)
    os.replace(tmp, out_path)


@tool(
    name="typeset_generate",
    description="把笔记/总结 Markdown 生成 PPTX 演示文稿（输出到 workspace/typeset/）",
    mode="typeset",
)
def typeset_generate(title: str, source: str = "notes", style: str = "auto") -> str:
    """把 notes/ 或 plans/ 下的一篇 Markdown 生成 PPTX，返回确认信息（含页数/主题/源字符数）。

    style 可选 auto / academic / deep / warm；auto 按源内容关键词猜主题；
    非法值返回结构化错误 dict。
    """
    if style not in ("auto",) + STYLE_NAMES:
        return _err(
            f"未知 style：{style!r}",
            "InvalidStyle",
            "style 可选值：auto / academic / deep / warm",
        )
    fname = _clean_title(title)
    workdir = _workdir()
    rel_src = os.path.join(source, fname)
    try:
        target_src = safe_read_path(rel_src, workdir)
    except ValueError as exc:
        return _denied_to_err(exc)
    if not os.path.isfile(target_src):
        return _err(
            f"{source} 下不存在：{os.path.join(source, fname)}",
            "NotFound",
            "用 notes/ 或 plans/ 下已有的文件名（可用 notes_read/notes_list 查看）",
        )
    with open(target_src, "r", encoding="utf-8") as f:
        content = f.read()
    theme_name = _resolve_theme(style, content)
    if theme_name is None:
        return _err(
            f"未知 style：{style!r}",
            "InvalidStyle",
            "style 可选值：auto / academic / deep / warm",
        )
    theme = THEMES[theme_name]
    slides = _build_slides(content, source, fname, target_src)
    out_name = fname[:-3] if fname.endswith(".md") else fname
    out_name += ".pptx"
    rel_out = os.path.join(TYPESET_DIR, out_name)
    try:
        out_path = safe_write_path(rel_out, workdir)
    except ValueError as exc:
        return _denied_to_err(exc)
    _save_pptx(slides, out_path, theme)
    return (
        f"已生成 {os.path.join(TYPESET_DIR, out_name)}"
        f"（{len(slides)} 页，style={theme_name}，{len(content)} 字符源内容）"
    )


def _extract_pdf_title(content: str) -> str | None:
    """取正文第一个一级标题作为 LaTeX \\title；无一级标题返回 None。"""
    for line in content.splitlines():
        m = re.match(r"^#\s+(.*)$", line.strip())
        if m:
            return m.group(1).strip()
    return None


def _escape_latex(text: str) -> str:
    """转义 LaTeX 特殊字符（& % # _ { } ~ ^；代码/数学上下文不经过这里）。"""
    for src, dst in (
        ("&", r"\&"),
        ("%", r"\%"),
        ("#", r"\#"),
        ("_", r"\_"),
        ("{", r"\{"),
        ("}", r"\}"),
        ("~", r"\textasciitilde{}"),
        ("^", r"\textasciicircum{}"),
    ):
        text = text.replace(src, dst)
    return text


def _apply_emphasis(text: str) -> str:
    """把剩余 **bold** / *italic* 记号转成 LaTeX 强调命令。"""
    text = re.sub(r"\*\*([^*]+)\*\*", r"\\textbf{\1}", text)
    text = re.sub(r"(?<!\*)\*([^*\s][^*]*?)\*(?!\*)", r"\\textit{\1}", text)
    return text


_INLINE_RE = re.compile(
    r"(\$\$.*?\$\$|\$[^$\n]*\$|`[^`]*`|!\[[^\]]*\]\([^)]*\)|\[[^\]]*\]\([^)]*\))",
    re.S,
)


def _latex_inline(text: str) -> str:
    """行内 Markdown → LaTeX：数学/代码/链接/图片先保护，其余转义后强调。"""
    out: list[str] = []
    pos = 0
    for m in _INLINE_RE.finditer(text):
        out.append(_escape_latex(text[pos : m.start()]))
        seg = m.group(1)
        if seg.startswith("$$") or seg.startswith("$"):
            out.append(seg)  # 数学公式原样透传
        elif seg.startswith("`"):
            out.append(r"\texttt{" + seg[1:-1] + "}")
        elif seg.startswith("!["):
            out.append("[图片略]")
        else:
            inner = re.match(r"\[([^\]]*)\]\(([^)]*)\)", seg)
            if inner:
                out.append(
                    r"\href{" + inner.group(2) + "}{"
                    + _latex_inline(inner.group(1)) + "}"
                )
            else:
                out.append(_escape_latex(seg))
        pos = m.end()
    out.append(_escape_latex(text[pos:]))
    return _apply_emphasis("".join(out))


def _latex_list(lines: list[str], i: int) -> tuple[str, int]:
    """列表块 → itemize/enumerate（4 空格缩进并入上级，生成嵌套列表）。

    返回 (LaTeX 片段, 下一行索引)。
    """
    items: list[tuple[int, str, bool]] = []
    n = len(lines)
    while i < n:
        raw = lines[i]
        if not raw.strip():
            break
        indent = len(raw) - len(raw.lstrip(" "))
        stripped = raw.strip()
        m_b = re.match(r"^[-*+]\s+(.*)$", stripped)
        m_o = re.match(r"^(\d+)[.)]\s+(.*)$", stripped)
        if m_b:
            items.append((indent, m_b.group(1), False))
            i += 1
            continue
        if m_o:
            items.append((indent, m_o.group(2), True))
            i += 1
            continue
        break

    def build(idx: int, indent: int) -> tuple[str, int]:
        env = "enumerate" if items[idx][2] else "itemize"
        out_lines = [f"\\begin{{{env}}}"]
        while idx < len(items) and items[idx][0] == indent:
            text = _latex_inline(items[idx][1])
            if idx + 1 < len(items) and items[idx + 1][0] > indent:
                inner, idx = build(idx + 1, items[idx + 1][0])
                out_lines.append(f"\\item {text}\n{inner}")
            else:
                out_lines.append(f"\\item {text}")
                idx += 1
        out_lines.append(f"\\end{{{env}}}")
        return "\n".join(out_lines), idx

    latex, _ = build(0, items[0][0])
    return latex, i


def _latex_table(rows: list[str]) -> str:
    """Markdown 表格 → tabular：首行表头 \\textbf，& 经行内转义为 \\&，表头上下 \\hline。"""
    parsed = _parse_table(rows)
    if parsed is None:
        return ""
    ncols = max(len(r) for r in parsed)
    header = parsed[0]
    body = parsed[1:]

    def cells(row: list[str]) -> str:
        padded = (row + [""] * ncols)[:ncols]
        return " & ".join(_latex_inline(c) for c in padded)

    lines = [
        "\\begin{table}[h]",
        "\\centering",
        f"\\begin{{tabular}}{{{'l' * ncols}}}",
        "\\hline",
        " & ".join(r"\textbf{" + _latex_inline(c) + "}" for c in header) + r" \\",
        "\\hline",
    ]
    for row in body:
        lines.append(cells(row) + r" \\")
    lines.extend(["\\hline", "\\end{tabular}", "\\end{table}"])
    return "\n".join(lines)


def _md_to_latex(md_text: str) -> str:
    """Markdown → LaTeX 正文。第一个 # 一级标题跳过（由 _build_latex 作 \\title）。"""
    lines = md_text.splitlines()
    out: list[str] = []
    i, n = 0, len(lines)
    first_h1 = True
    while i < n:
        raw = lines[i].rstrip()
        stripped = raw.strip()
        if not stripped:
            i += 1
            continue
        if stripped.startswith("```"):
            code = []
            i += 1
            while i < n and not lines[i].strip().startswith("```"):
                code.append(lines[i])
                i += 1
            i += 1  # 跳过闭合 fence
            out.append("\\begin{verbatim}\n" + "\n".join(code) + "\n\\end{verbatim}")
            continue
        m_head = re.match(r"^(#{1,6})\s+(.*)$", stripped)
        if m_head:
            level = len(m_head.group(1))
            text = _latex_inline(m_head.group(2).strip())
            if level == 1 and first_h1:
                first_h1 = False
                i += 1
                continue
            if level <= 2:
                out.append(f"\\section{{{text}}}")
            elif level == 3:
                out.append(f"\\subsection{{{text}}}")
            else:
                out.append(f"\\subsubsection{{{text}}}")
            i += 1
            continue
        if stripped.startswith("|"):
            rows = [stripped]
            i += 1
            while i < n and lines[i].strip().startswith("|"):
                rows.append(lines[i].strip())
                i += 1
            out.append(_latex_table(rows))
            continue
        if stripped.startswith(">"):
            quote = _latex_inline(stripped.lstrip(">").strip())
            out.append(f"\\begin{{quote}}{quote}\\end{{quote}}")
            i += 1
            continue
        if re.match(r"^!\[[^\]]*\]\([^)]*\)\s*$", stripped):
            out.append(f"% 图片略：{stripped}")
            i += 1
            continue
        if re.match(r"^[-*+]\s+", stripped) or re.match(r"^\d+[.)]\s+", stripped):
            latex, i = _latex_list(lines, i)
            out.append(latex)
            continue
        para = [raw]
        i += 1
        while i < n:
            nxt = lines[i].rstrip()
            ns = nxt.strip()
            if not ns:
                break
            if re.match(r"^(#{1,6})\s+", ns) or ns.startswith("```") \
                    or ns.startswith("|") or ns.startswith(">") \
                    or re.match(r"^!\[[^\]]*\]\([^)]*\)\s*$", ns) \
                    or re.match(r"^[-*+]\s+", ns) or re.match(r"^\d+[.)]\s+", ns):
                break
            para.append(nxt)
            i += 1
        out.append(_latex_inline(" ".join(p.strip() for p in para)))
    return "\n\n".join(out)


def _build_latex(title: str, body: str, source_note: str) -> str:
    """拼 LaTeX 全文（ctexart + geometry + amsmath + hyperref + booktabs）。

    source_note 预留参数（当前模板写死，未使用）。
    """
    return (
        "\\documentclass[11pt]{ctexart}\n"
        "\\usepackage[margin=2.5cm]{geometry}\n"
        "\\usepackage{amsmath,amssymb}\n"
        "\\usepackage{hyperref}\n"
        "\\usepackage{booktabs}\n"
        "\\hypersetup{colorlinks=true,linkcolor=blue!60!black,citecolor=blue!60!black,urlcolor=blue!60!black}\n"
        f"\\title{{{_escape_latex(title)}}}\n"
        "\\author{PhySc-agent}\n"
        f"\\date{{{datetime.now().strftime('%Y-%m-%d')}}}\n"
        "\\begin{document}\n"
        "\\maketitle\n"
        f"{body}\n"
        "\\end{document}\n"
    )


def _compile_xelatex(tex_path: str, out_dir: str) -> tuple[int, str]:
    """xelatex 编译：返回 (returncode, log 尾部 800 字符)；超时抛 RuntimeError。"""
    try:
        proc = subprocess.run(
            [
                "xelatex",
                "-interaction=nonstopmode",
                "-halt-on-error",
                "-output-directory",
                out_dir,
                tex_path,
            ],
            capture_output=True,
            text=True,
            timeout=120,
        )
    except subprocess.TimeoutExpired:
        raise RuntimeError("编译超时")
    return proc.returncode, proc.stdout[-800:]


def _pdf_pages(pdf_path: str) -> int | None:
    """用 pymupdf 打开 PDF 取页数；打开失败返回 None。"""
    try:
        with _fitz.open(pdf_path) as doc:
            return doc.page_count
    except Exception:
        return None


@tool(
    name="typeset_pdf",
    description="把 notes/ 或 plans/ 下的 Markdown 笔记生成 PDF 学术文档：md → LaTeX（ctexart 中文排版 + amsmath 数学公式 + verbatim 代码块）→ xelatex 编译 → 输出 workspace/typeset/<标题>.pdf（同时保留 .tex 源文件）。支持表格与列表。返回 PDF 路径 + 页数（若可得）",
    mode={"typeset"},
)
def typeset_pdf(title: str, source: str = "notes", style: str = "academic") -> dict:
    """把 notes/ 或 plans/ 下的一篇 Markdown 生成 PDF（ctexart 学术白底），保留 .tex 源文件。

    style 仅作兼容保留：PDF 一律 academic 白底学术风，非 academic 值忽略并提示。
    """
    if style != "academic":
        style_note = f"（style={style!r} 已忽略，PDF 一律 academic 白底学术版式）"
    else:
        style_note = ""
    fname = _clean_title(title)
    workdir = _workdir()
    rel_src = os.path.join(source, fname)
    try:
        target_src = safe_read_path(rel_src, workdir)
    except ValueError as exc:
        return _denied_to_err(exc)
    if not os.path.isfile(target_src):
        return _err(
            f"{source} 下不存在：{os.path.join(source, fname)}",
            "NotFound",
            "用 notes/ 或 plans/ 下已有的文件名（可用 notes_read/notes_list 查看）",
        )
    with open(target_src, "r", encoding="utf-8") as f:
        content = f.read()
    base = fname[:-3] if fname.endswith(".md") else fname
    if not base:
        base = "untitled"
    title_latex = _extract_pdf_title(content) or base
    latex = _build_latex(title_latex, _md_to_latex(content), f"{source}/{fname}")
    out_dir = os.path.join(workdir, TYPESET_DIR)
    os.makedirs(out_dir, exist_ok=True)
    tex_path = os.path.join(out_dir, base + ".tex")
    with open(tex_path, "w", encoding="utf-8") as f:
        f.write(latex)
    try:
        returncode, log_tail = _compile_xelatex(tex_path, out_dir)
    except RuntimeError as exc:
        return _err(str(exc), "CompileTimeout", "重试，或确认系统已安装 xelatex（texlive-full）")
    if returncode != 0:
        return _err(
            f"PDF 编译失败：{log_tail[-200:]}",
            "CompileError",
            "检查笔记中的特殊字符或表格格式",
        )
    rel = f"{TYPESET_DIR}/{base}"
    return {
        "pdf_file": f"workspace/{rel}.pdf",
        "tex_file": f"workspace/{rel}.tex",
        "pages": _pdf_pages(os.path.join(out_dir, base + ".pdf")),
        "note": f"PDF 已生成：workspace/{rel}.pdf{style_note}",
    }
